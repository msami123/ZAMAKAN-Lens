import pandas as pd
import base64
import streamlit as st
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

def load_data(file_path):
    if file_path.name.endswith('.csv'):
        return pd.read_csv(file_path)
    else:
        return pd.read_excel(file_path)

def save_figure_to_image(fig, width=1200, height=675):
    """Save the Plotly figure as an image with specified size."""
    buffer = BytesIO()
    fig.update_layout(width=width, height=height)  # Adjust size for widescreen
    fig.write_image(buffer, format='png', engine='kaleido')
    buffer.seek(0)
    return buffer

def create_pptx(figures):
    """Create a PowerPoint presentation with properly sized wide figures."""
    ppt = Presentation()
    slide_width, slide_height = Inches(13.33), Inches(7.5)  # Widescreen 16:9 slide size

    for title, fig in figures.items():
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Add a blank slide
        img_stream = save_figure_to_image(fig)

        # Set image size to use the full slide width while maintaining aspect ratio
        img_width = slide_width
        img_height = img_width * 9 / 16  # Maintain 16:9 aspect ratio

        # Calculate centered position
        left, top = (slide_width - img_width) / 2, (slide_height - img_height) / 2

        # Add the image to the slide
        slide.shapes.add_picture(img_stream, left, top, width=img_width, height=img_height)

    # Save the presentation to a BytesIO buffer
    ppt_buffer = BytesIO()
    ppt.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

st.set_page_config(layout="wide")
st.title('Your Audience Insights Begins Here.')

# Sidebar for file upload and settings
with st.sidebar:
    try:
        with open(r"/Users/anody/Downloads/iAMAi-Viz-logo-03.png", "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode("utf-8")
        img = f"<img src='data:image/png;base64,{encoded_string}' alt='App Logo' style='width:auto;height:auto;'>"
        st.markdown(img, unsafe_allow_html=True)
    except FileNotFoundError:
        st.error("Logo image not found. Please ensure the path is correct.")

    uploaded_file = st.file_uploader("Choose a data file (CSV or XLSX)", type=['csv', 'xlsx'])
    color = st.color_picker('Pick a color for the bars/lines', '#468294', key="color")
    bgcolor = st.color_picker('Pick a background color', '#101116', key="bgcolor")
    title_size = st.slider("Title Font Size", 10, 50, 24, key="title_size")
    axis_label_size = st.slider("Axis Label Font Size", 10, 30, 18, key="axis_label_size")
    legend_size = st.slider("Legend Font Size", 10, 30, 14, key="legend_size")

    with st.expander("About IAMAI VIZ !!", expanded=False):
        st.markdown("""
        **IAMAI VIZ: Data Visualization Tool**

        This interactive tool helps users visualize data efficiently. Users can generate insights and explore data through various visualizations.

        **developed by M&M**
        """)

if uploaded_file is not None:
    df = load_data(uploaded_file)

    st.subheader('Data Preview')
    with st.expander('Data Preview'):
        st.dataframe(df)

    st.subheader('Data Visualization')
    figures = {}  # Store all figures for PowerPoint export

    with st.expander('Visualize Data'):
        if 'Short Label Question' in df.columns:
            short_label_questions = df['Short Label Question'].unique()
            for question in short_label_questions:
                st.subheader(question)
                visualization_type = st.selectbox(
                    f'Select the visualization type for {question}',
                    ['Bar Chart', 'Line Chart', 'Pie Chart'],
                    key=question
                )
                question_data = df[df['Short Label Question'] == question]

                if visualization_type == 'Bar Chart':
                    fig = px.bar(
                        question_data, x='Attributes', y='Audience %',
                        title=question, color='Attributes', text='Audience %'
                    )
                    fig.update_traces(marker=dict(color=color))
                    fig.update_layout(
                        xaxis_title='Attributes', yaxis_title='Audience %',
                        height=800, title_x=0.5, plot_bgcolor=bgcolor,
                        title_font_size=title_size,
                        xaxis_title_font_size=axis_label_size,
                        yaxis_title_font_size=axis_label_size,
                        legend_title_font_size=legend_size
                    )
                    fig.update_xaxes(tickangle=90)
                    fig.update_traces(texttemplate='%{text:.2f}', textposition='outside')
                    st.plotly_chart(fig)
                    figures[f"{question}_bar"] = fig

                elif visualization_type == 'Line Chart':
                    fig = px.line(question_data, x='Attributes', y='Audience %', title=question)
                    fig.update_traces(line=dict(color=color))
                    fig.update_layout(
                        xaxis_title='Attributes', yaxis_title='Audience %',
                        height=800, title_x=0.5, plot_bgcolor=bgcolor,
                        title_font_size=title_size,
                        xaxis_title_font_size=axis_label_size,
                        yaxis_title_font_size=axis_label_size,
                        legend_title_font_size=legend_size
                    )
                    fig.update_xaxes(tickangle=90)
                    st.plotly_chart(fig)
                    figures[f"{question}_line"] = fig

                elif visualization_type == 'Pie Chart':
                    fig = px.pie(
                        question_data, values='Audience %', names='Attributes', title=question,
                        color_discrete_sequence=px.colors.qualitative.Set3  # Ensure multiple colors
                    )
                    fig.update_layout(
                        height=600, title_x=0.5, plot_bgcolor=bgcolor,
                        title_font_size=title_size,
                        legend_title_font_size=legend_size
                    )
                    st.plotly_chart(fig)
                    figures[f"{question}_pie"] = fig
        else:
            st.write("The dataset does not contain a column named 'Short Label Question'.")

    if figures:
        ppt_buffer = create_pptx(figures)
        st.download_button(
            label="Download All Graphs as PowerPoint",
            data=ppt_buffer,
            file_name="graphs.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
